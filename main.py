import google.generativeai as genai
import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


# Gemini Setup
os.environ["GEMINI_API_KEY"] = "*****"
genai.configure(api_key=os.environ["GEMINI_API_KEY"])
model = genai.GenerativeModel(model_name="gemini-1.5-flash")

system_instruction = (
    "You are a study buddy app, i give you text/info about something and you create a quizes, flash cards and summaries based on the information i give you"
)

chat = model.start_chat(history=[{"role": "user", "parts": [system_instruction]}])

recall = False
user_input = ""

user_initial_input = input(
    "Choose input type: normal text, pdf file, word file, pptx file, or youtube link: "
).lower()

match user_initial_input:
    case "normal text":
        user_input = input("please enter the text: ")

    case "pdf file":
        pdf_path = input("please give file path to pdf file: ")

        with fitz.open(pdf_path) as doc:
            text = ""
            for page in doc:
                text += page.get_text()
        user_input = text

    case "word file":
        word_path = input("please give file path to pdf file: ")

        doc = Document(word_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        user_input = text

    case "pptx file":
        pptx_path = input("please give file path to pdf file: ")
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        user_input = text

    case "youtube link": # not currently working
        youtube_link = input("please enter youtube link: ")
        user_input = "I am going to give you and youtube link: " + youtube_link


user_choice = input("Select: quiz, flash cards, summary, or recall - ").lower()

match user_choice:
    case "quiz":
        recall = False
        initial_command = "Create a quiz based on the following text/youtube link in JSON format: "

    case "flash cards":
        recall = False
        initial_command = "Create flash cards based on the following text/youtube link in JSON format: "

    case "summary":
        recall = False
        initial_command = "Create a summary of the following text/youtube link in JSON format: "

    case "recall what i have learned":
        recall = True
        initial_command = ("I am going to give you two texts, the first is going to be the lesson, "
                           "and the second on is what i have learned. "
                           "i need you to compare the second one to the first one and give a percentage of how much of the lesson i have rememmbered: ")
        recall_what_i_have_learned = input("Enter what have you learned - ")


#user_input = ("The solar system consists of the Sun and eight planets: Mercury, Venus, Earth, Mars, Jupiter, Saturn, Uranus, and Neptune. "
              # "Jupiter is the largest planet, while Mercury is the smallest. "
              # "The Sun is a star, not a planet, and it provides energy through nuclear fusion. "
              # "Earth is the only planet known to support life.")

if recall:
    send_to_model = initial_command + user_input + ", and the second one: " + recall_what_i_have_learned
else:
    send_to_model = initial_command + user_input

result = chat.send_message({"parts": [send_to_model]})
response_text = result.text
print(f"Quiz bot: {response_text}")
