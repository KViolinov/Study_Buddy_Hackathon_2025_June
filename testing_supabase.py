# from supabase import create_client, Client
#
# # url = os.getenv("SUPABASE_URL")
# # key = os.getenv("SUPABASE_KEY")
#
# url = "https://qjbwcxqxtgftawrqxuec.supabase.co"
# key = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InFqYndjeHF4dGdmdGF3cnF4dWVjIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NTA1MTg5NDcsImV4cCI6MjA2NjA5NDk0N30.DWwk1lldHiuZwqocTFMQD3w-cNt8geJaeq1ZeRT84Go"
#
# supabase: Client = create_client(url, key)
#
# # Define the new record as a dictionary
# new_input = {
#     "user_id": "1",
#     "time_of_sending": "2023-10-01T12:00:00Z",
#     "type": "summary",
#     "input_source_type": "normal text",
#     "input_text": "This is a sample input text for testing purposes.",
# }
#
# # Insert into the Inputs table
# response = supabase.table("inputs").insert(new_input).execute()
#
# print("Inserted:", response.data)






# import time
#
# from supabase import create_client, Client
#
# # url = os.getenv("SUPABASE_URL")
# # key = os.getenv("SUPABASE_KEY")
#
# url = "https://qjbwcxqxtgftawrqxuec.supabase.co"
# key = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InFqYndjeHF4dGdmdGF3cnF4dWVjIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NTA1MTg5NDcsImV4cCI6MjA2NjA5NDk0N30.DWwk1lldHiuZwqocTFMQD3w-cNt8geJaeq1ZeRT84Go"
#
# supabase: Client = create_client(url, key)
#
# while True:
#     # Fetch the latest entry based on time_of_sending
#     response = (
#         supabase
#         .table("inputs")
#         .select("*")
#         .order("time_of_sending", desc=True)
#         .limit(1)
#         .execute()
#     )
#
#     # Extract and print the single latest row
#     if response.data:
#         latest = response.data[0]
#         print("Latest entry:", latest)
#     else:
#         print("No entries found.")
#
#     user_id = latest["user_id"]
#     time_of_sending = latest["time_of_sending"]
#     input_type = latest["type"]
#     input_source_type = latest["input_source_type"]
#     input_text = latest["input_text"]
#
#     print("sleeping for 15 seconds to avoid rate limiting")
#     time.sleep(15)






import time
import google.generativeai as genai
from datetime import datetime, timezone
import os
import json
import fitz
from docx import Document
from pptx import Presentation

from supabase import create_client, Client

from dotenv import load_dotenv
load_dotenv()

url = os.getenv("SUPABASE_URL")
key = os.getenv("SUPABASE_KEY")

# url = "https://qjbwcxqxtgftawrqxuec.supabase.co"
# key = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InFqYndjeHF4dGdmdGF3cnF4dWVjIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NTA1MTg5NDcsImV4cCI6MjA2NjA5NDk0N30.DWwk1lldHiuZwqocTFMQD3w-cNt8geJaeq1ZeRT84Go"

supabase: Client = create_client(url, key)

# Gemini Setup
os.environ["GEMINI_API_KEY"] = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=os.environ["GEMINI_API_KEY"])
model = genai.GenerativeModel(model_name="gemini-1.5-flash")

system_instruction = (
    "You are a study buddy app, i give you full text/lesson about something and you create a quizes, "
    "flash cards and/or summaries based on the information i give you"
)

chat = model.start_chat(history=[{"role": "user", "parts": [system_instruction]}])


# def summarize_youtube_video(youtube_link):
#     try:
#         # Extract video ID from the link
#         video_id = youtube_link.split("watch?v=")[1]
#         if "&" in video_id:
#             video_id = video_id.split("&")[0]
#
#         # Fetch the transcript
#         try:
#             from youtube_transcript_api import YouTubeTranscriptApi
#             transcript = YouTubeTranscriptApi.get_transcript(video_id)
#             text = ' '.join([entry['text'] for entry in transcript])
#         except Exception as e:
#             return f"Error fetching transcript: {e}"
#
#         if not text:
#             return "Could not retrieve transcript for the video."
#
#         return text  # Return the transcript
#     except Exception as e:
#         return f"An error occurred: {e}"

recall = False
user_input = ""

last_time_of_sending = None

while True:
    # Fetch the latest entry based on time_of_sending
    response = (
        supabase
        .table("inputs")
        .select("*")
        .order("time_of_sending", desc=True)
        .limit(1)
        .execute()
    )

    if response.data:
        latest = response.data[0]
        if latest["time_of_sending"] == last_time_of_sending:
            print("No new entry. Sleeping...")
            time.sleep(15)
            continue  # Skip processing if not new

        last_time_of_sending = latest["time_of_sending"]
        print("Latest entry:", latest)
    else:
        print("No entries found.")
        time.sleep(15)
        continue

    user_id = latest["user_id"]
    input_id = latest["input_id"]
    time_of_sending = latest["time_of_sending"]
    input_type = latest["type"]
    input_source_type = latest["input_source_type"]
    input_text = latest["input_text"]

    match input_source_type:
        case "normal text":
            user_input = input_text

        # case "pdf file":
        #     pdf_path = input("please give file path to pdf file: ")
        #
        #     with fitz.open(pdf_path) as doc:
        #         text = ""
        #         for page in doc:
        #             text += page.get_text()
        #     user_input = text
        #
        # case "word file":
        #     word_path = input("please give file path to pdf file: ")
        #
        #     doc = Document(word_path)
        #     text = "\n".join([para.text for para in doc.paragraphs])
        #     user_input = text
        #
        # case "pptx file":
        #     pptx_path = input("please give file path to pdf file: ")
        #     prs = Presentation(pptx_path)
        #     text = ""
        #     for slide in prs.slides:
        #         for shape in slide.shapes:
        #             if hasattr(shape, "text"):
        #                 text += shape.text + "\n"
        #     user_input = text
        #
        # case "youtube link":  # now working
        #     youtube_link = input("please enter youtube link: ")
        #     user_input = summarize_youtube_video(youtube_link)
        #     if "Error" in user_input:
        #         print(user_input)
        #         user_input = ""  # Reset user_input if there was an error

    match input_type:
        case "quiz":
            recall = False
            initial_command = ("Create a quiz in JSON format: following"
                               "this strict model format-> your json has to have similar format like this one about quiz:"
                               "public class Quiz {"
                               "public string Question { get; set; }"
                               "public string CorrectAnswer { get; set; }"
                               "public List<string> Options { get; set; }")

        case "flash cards":
            recall = False
            initial_command = "Create flash cards based on the following text in JSON format: "

        case "summary":
            recall = False
            initial_command = ("Create a summarycase summary")
            recall = False
            initial_command = ("Create a summary of the following text/youtube link in JSON format:"
            "The returned json should be in this format:"
            "public class Summary {"
                               
            "public string summary { get; set; }"
            "public string[] keywords { get; set; } in JSON format:"
            "The returned json should be in this format:"
            "public class Summary {"
            "public string summary { get; set; }"
            "public string[] keywords { get; set; }")

        case "recall":
            recall = True
            initial_command = ("I am going to give you two texts, the first is going to be the lesson, "
                               "and the second on is what i have learned. "
                               "i need you to compare the second one to the first one "
                               "and give a percentage of how much of the lesson i have remembered with tips to improve, but without summary or flashcards"
                               "All of this needs to be in JSON format: ")
            recall_what_i_have_learned = input("Enter what have you learned - ")

    if recall:
        send_to_model = initial_command + user_input + ", and the second one: " + recall_what_i_have_learned
    else:
        send_to_model = initial_command + user_input

    result = chat.send_message({"parts": [send_to_model]})
    response_text = result.text
    print(f"Quiz bot: {response_text}")

    print("Inserting into the database...")

    now_iso = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

    # Define the new record as a dictionary
    new_output = {
        "user_id": user_id,
        "time_of_sending": now_iso,
        "type": input_type,
        "output_text": response_text,
        'input_id': input_id,
    }

    # Insert into the Inputs table
    response = supabase.table("outputs").insert(new_output).execute()

    print("Inserted:", response.data)

    print("sleeping for 15 seconds to avoid rate limiting")
    time.sleep(15)

    # Po vibe kodnato nqma na kude brat
